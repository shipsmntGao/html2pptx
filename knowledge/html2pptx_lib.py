# html2pptx_lib.py
"""
html2pptx_lib.py — Lightweight HTML→PPTX converter (v2, images & tables supported).

Project: プレゼン資料作成v7
Author: GPT-5 Thinking
License: MIT

What’s new in v2
----------------
- <img>, <figure><img><figcaption> supported in .content blocks
  - Only local file paths or data: URIs are supported. External HTTP(S) URLs are ignored.
- <table> supported (with optional <thead> and <tbody>)
  - Header row (from <thead> or first <tr> with <th>) rendered in bold with accent background
  - Equal column widths; zebra striping optional (disabled by default for cleanliness)
- Mixed content flow inside .content rendered top→bottom in a simple flow layout:
  paragraphs / headings / lists / images / tables.
  Heights are approximated for text to avoid overlap; keep paragraphs concise and prefer
  one visual (image or table) per .content to ensure clean layouts.

Supported HTML subset (strict)
------------------------------
- Article wrapper: <article class="slide [title-slide]?">…</article>
- Header: <header class="header"><div class="eyebrow">…</div><h2 class="title">…</h2></header>
- Content: <section class="content [two-col]?">…</section>
  Children may include:
    <h1|h2|h3>, <p>, <div>, <ul><li>…</li></ul>, <img src="...">,
    <figure><img src="..."><figcaption>…</figcaption></figure>,
    <table> with <thead>/<tbody>/<tr>/<th>/<td>
- Footer: <footer class="footer"><div>left</div><div>page markers</div></footer>

Dependencies expected in Code Interpreter:
  - python-pptx
  - beautifulsoup4

Example:
    from html2pptx_lib import convert_html_to_pptx
    convert_html_to_pptx("draft.html", "output.pptx")
"""
from __future__ import annotations

import re
import json
import base64
import uuid
import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Optional, Tuple, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

try:
    from bs4 import BeautifulSoup  # type: ignore
except Exception as e:  # pragma: no cover
    raise RuntimeError("beautifulsoup4 is required (pip install beautifulsoup4).") from e


# ---------- Utilities ----------

def hex_to_rgb(hex_color: str) -> Tuple[int,int,int]:
    hex_color = hex_color.strip()
    if hex_color.startswith("#"):
        hex_color = hex_color[1:]
    if len(hex_color) == 3:
        hex_color = "".join([c*2 for c in hex_color])
    if len(hex_color) != 6:
        return (0,0,0)
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return (r,g,b)

def rgba_to_rgb(rgba: str) -> Tuple[int,int,int]:
    # Accepts 'rgba(r,g,b,a)' or 'rgb(r,g,b)' — returns (r,g,b), ignoring alpha.
    nums = re.findall(r"[\d.]+", rgba)
    if len(nums) >= 3:
        r, g, b = [int(float(nums[i])) for i in range(3)]
        return (r, g, b)
    return (0,0,0)

def parse_css_vars(style_text: str) -> Dict[str, str]:
    """
    Extracts CSS variables from the :root scope in the <style> element.
    Returns a mapping like {'--brand': '#0f3d68', '--fs-title': 'clamp(24px, 3.8vmin, 48px)', ...}
    """
    m = re.search(r":root\s*\{([^}]*)\}", style_text, re.DOTALL)
    if not m:
        return {}
    block = m.group(1)
    vars_map: Dict[str, str] = {}
    for line in block.splitlines():
        if ":" in line and "--" in line:
            k_v = line.strip().rstrip(";")
            if not k_v:
                continue
            parts = k_v.split(":", 1)
            if len(parts) != 2:
                continue
            key = parts[0].strip()
            val = parts[1].strip()
            vars_map[key] = val
    return vars_map

def eval_clamp_px(expr: str, vw: int = 1920, vh: int = 1080) -> float:
    """
    Evaluate a CSS clamp() or raw px value at baseline viewport (1920×1080).
    Supports forms:
      - 'clamp(24px, 3.8vmin, 48px)'
      - '32px'
      - '3.2vmin'
    Returns a pixel value (float).
    """
    expr = expr.strip()
    def vmin_to_px(vmin_str: str) -> float:
        vmin = min(vw, vh) / 100.0
        num = float(re.findall(r"[-+]?\d*\.?\d+", vmin_str)[0])
        return num * vmin

    def to_px(val: str) -> float:
        val = val.strip()
        if val.endswith("px"):
            return float(val[:-2])
        if val.endswith("vmin"):
            return vmin_to_px(val)
        try:
            return float(val)
        except:
            return 0.0

    if expr.startswith("clamp(") and expr.endswith(")"):
        inner = expr[len("clamp("):-1]
        parts = [p.strip() for p in inner.split(",")]
        if len(parts) == 3:
            a, b, c = parts
            mid = to_px(b)
            low = to_px(a)
            high = to_px(c)
            return max(low, min(mid, high))
    return to_px(expr)

def px_to_pt(px: float) -> float:
    return px * 0.75


@dataclass
class DesignTokens:
    brand: Tuple[int,int,int]
    accent: Tuple[int,int,int]
    bg: Tuple[int,int,int]
    ink: Tuple[int,int,int]
    muted: Tuple[int,int,int]
    line: Tuple[int,int,int]
    font_sans: str
    fs_title_xl_pt: float
    fs_title_pt: float
    fs_body_pt: float
    fs_small_pt: float

    @classmethod
    def from_css_vars(cls, varmap: Dict[str,str]) -> "DesignTokens":
        def color(key: str, default_hex: str="#000000") -> Tuple[int,int,int]:
            raw = varmap.get(key, default_hex)
            if raw.startswith("rgba(") or raw.startswith("rgb("):
                return rgba_to_rgb(raw)
            if raw.startswith("#"):
                return hex_to_rgb(raw)
            return hex_to_rgb(default_hex)

        def font(key: str, default: str="Noto Sans JP"):
            raw = varmap.get(key, default)
            fam = raw.split(",")[0].strip().strip('"').strip("'")
            return fam or default

        def fs(key: str, default_px: float) -> float:
            raw = varmap.get(key, f"{default_px}px")
            px = eval_clamp_px(raw)
            return px_to_pt(px)

        return cls(
            brand=color("--brand", "#0f3d68"),
            accent=color("--accent", "#2563eb"),
            bg=color("--bg", "#ffffff"),
            ink=color("--ink", "#0b1220"),
            muted=color("--muted", "#6b7280"),
            line=color("--line", "#93a1b0"),
            font_sans=font("--font-sans", "Noto Sans JP"),
            fs_title_xl_pt=fs("--fs-title-xl", 48),
            fs_title_pt=fs("--fs-title", 36),
            fs_body_pt=fs("--fs-body", 18),
            fs_small_pt=fs("--fs-small", 12),
        )


# ---------- Core converter ----------

class HTMLtoPPTXConverter:
    def __init__(self, *, page_width_in=13.3333, page_height_in=7.5):
        self.page_width_in = page_width_in
        self.page_height_in = page_height_in
        self._base_dir: Optional[Path] = None
        self.warnings: List[str] = []

    def _rgb(self, prs_color: Tuple[int,int,int]) -> RGBColor:
        r,g,b = prs_color
        return RGBColor(r,g,b)

    def _set_font(self, run, *, name: str, size_pt: float, color: Tuple[int,int,int], bold: bool=False):
        run.font.name = name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color)

    def _add_footer(self, slide, tokens: DesignTokens, left_text: str, right_text: str):
        left_margin = Inches(0.6)
        right_margin = Inches(0.6)
        footer_top = Inches(self.page_height_in - 0.7)
        width = Inches(self.page_width_in - 1.2)
        height = Inches(0.5)

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_margin, footer_top - Inches(0.12), width, Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(tokens.line)
        line.line.fill.background()

        left_tb = slide.shapes.add_textbox(left_margin, footer_top, Inches(6.0), height).text_frame
        left_tb.clear()
        p = left_tb.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        p.line_spacing = 1.5
        run = p.add_run()
        self._set_font(run, name=tokens.font_sans, size_pt=tokens.fs_small_pt, color=tokens.muted, bold=False)
        run.text = left_text

        right_tb = slide.shapes.add_textbox(Inches(self.page_width_in - 0.6 - 2.0), footer_top, Inches(2.0), height).text_frame
        right_tb.clear()
        p2 = right_tb.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        p2.space_after = Pt(0)
        p2.line_spacing = 1.5
        run2 = p2.add_run()
        self._set_font(run2, name=tokens.font_sans, size_pt=tokens.fs_small_pt, color=tokens.muted, bold=False)
        run2.text = right_text

    def _add_header(self, slide, tokens: DesignTokens, eyebrow: Optional[str], title: Optional[str]):
        left_margin = Inches(0.6)
        top = Inches(0.4)
        width = Inches(self.page_width_in - 1.2)
        height = Inches(1.4)

        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_margin, top + height, width, Inches(0.01))
        rule.fill.solid()
        rule.fill.fore_color.rgb = self._rgb(tokens.line)
        rule.line.fill.background()

        if eyebrow:
            tb1 = slide.shapes.add_textbox(left_margin, top, width, Inches(0.4)).text_frame
            tb1.clear()
            p1 = tb1.paragraphs[0]
            p1.alignment = PP_ALIGN.LEFT
            p1.space_after = Pt(0)
            p1.line_spacing = 1.5
            r1 = p1.add_run()
            self._set_font(r1, name=tokens.font_sans, size_pt=tokens.fs_small_pt, color=tokens.muted, bold=False)
            r1.text = eyebrow

        if title:
            tb2 = slide.shapes.add_textbox(left_margin, top + Inches(0.35), width, Inches(1.0)).text_frame
            tb2.clear()
            p2 = tb2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(0)
            p2.line_spacing = 1.5
            r2 = p2.add_run()
            self._set_font(r2, name=tokens.font_sans, size_pt=tokens.fs_title_pt, color=tokens.ink, bold=True)
            r2.text = title

    def _add_title_slide(self, prs: Presentation, tokens: DesignTokens, kicker: str, title: str, subtitle: str, footer_left: str, page_text: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(self.page_width_in-2.0), Inches(3.0)).text_frame
        tb.clear()

        p0 = tb.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(4)
        p0.line_spacing = 1.5
        r0 = p0.add_run()
        self._set_font(r0, name=tokens.font_sans, size_pt=tokens.fs_small_pt, color=tokens.muted, bold=False)
        r0.text = kicker

        p1 = tb.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(6)
        p1.line_spacing = 1.2
        r1 = p1.add_run()
        self._set_font(r1, name=tokens.font_sans, size_pt=tokens.fs_title_xl_pt, color=tokens.ink, bold=True)
        r1.text = title

        p2 = tb.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(0)
        p2.line_spacing = 1.5
        r2 = p2.add_run()
        self._set_font(r2, name=tokens.font_sans, size_pt=tokens.fs_body_pt, color=tokens.muted, bold=False)
        r2.text = subtitle

        self._add_footer(slide, tokens, footer_left, page_text)
        return slide

    # ---------- Images & Tables helpers ----------

    def _decode_data_uri(self, data_uri: str) -> Optional[Path]:
        try:
            # data:[<mediatype>][;base64],<data>
            header, b64 = data_uri.split(',', 1)
            mime = "application/octet-stream"
            if ';' in header:
                parts = header.split(';')
                if parts[0].startswith('data:'):
                    mime = parts[0][5:] or mime
            ext = mimetypes.guess_extension(mime) or ".bin"
            tmp = Path("/mnt/data") / f"img_{uuid.uuid4().hex}{ext}"
            tmp.write_bytes(base64.b64decode(b64))
            return tmp
        except Exception as e:
            self.warnings.append(f"Failed to decode data URI image: {e}")
            return None

    def _resolve_image_path(self, src: str) -> Optional[Path]:
        if not src:
            return None
        s = src.strip()
        if s.startswith("data:"):
            return self._decode_data_uri(s)
        if re.match(r"^https?://", s, flags=re.I):
            self.warnings.append(f"Skip remote image (unsupported in this environment): {s}")
            return None
        p = Path(s)
        if not p.is_absolute():
            base = self._base_dir or Path(".")
            p = base / s
        if p.exists():
            return p
        self.warnings.append(f"Image not found: {p}")
        return None

    def _line_height_in(self, tokens: DesignTokens, *, body: bool=True) -> float:
        pt = tokens.fs_body_pt if body else tokens.fs_small_pt
        return (pt * 1.5) / 72.0

    def _add_text_block(self, slide, tokens: DesignTokens, lines: List[str], x_in: float, y_in: float, w_in: float, *, bold: bool=False, small: bool=False) -> float:
        """Add a stack of lines as a textbox. Returns height used (inches)."""
        est_lines = max(1, len(lines))
        line_h = self._line_height_in(tokens, body=not small)
        block_h = est_lines * line_h + 0.05
        tf = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(block_h)).text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)
            p.line_spacing = 1.5
            run = p.add_run()
            if small:
                self._set_font(run, name=tokens.font_sans, size_pt=tokens.fs_small_pt, color=tokens.muted, bold=bold)
            else:
                self._set_font(run, name=tokens.font_sans, size_pt=tokens.fs_body_pt, color=tokens.ink, bold=bold)
            run.text = line
        return block_h

    def _add_list_block(self, slide, tokens: DesignTokens, items: List[str], x_in: float, y_in: float, w_in: float) -> float:
        est_lines = max(1, len(items))
        line_h = self._line_height_in(tokens, body=True)
        block_h = est_lines * line_h + 0.1
        tf = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(block_h)).text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)
            p.line_spacing = 1.5
            run = p.add_run()
            self._set_font(run, name=tokens.font_sans, size_pt=tokens.fs_body_pt, color=tokens.ink, bold=False)
            run.text = f"・ {item}"
        return block_h

    def _add_image_block(self, slide, tokens: DesignTokens, src: str, x_in: float, y_in: float, w_in: float) -> float:
        p = self._resolve_image_path(src)
        if not p:
            return 0.0
        pic = slide.shapes.add_picture(str(p), Inches(x_in), Inches(y_in), width=Inches(w_in))
        used_h = float(pic.height) / 914400.0  # emu→inch
        return used_h

    def _add_table_block(self, slide, tokens: DesignTokens, table_dom, x_in: float, y_in: float, w_in: float) -> float:
        # Parse headers
        headers: List[str] = []
        body_rows: List[List[str]] = []
        thead = table_dom.find("thead")
        tbody = table_dom.find("tbody")
        rows = []
        if thead:
            rows.extend(thead.find_all("tr", recursive=False))
        if tbody:
            rows.extend(tbody.find_all("tr", recursive=False))
        if not rows:
            rows = table_dom.find_all("tr", recursive=False)

        # Detect header row
        if rows and rows[0].find_all("th"):
            headers = [c.get_text(" ", strip=True) for c in rows[0].find_all(["th","td"], recursive=False)]
            data_trs = rows[1:]
        else:
            data_trs = rows

        for tr in data_trs:
            cells = [c.get_text(" ", strip=True) for c in tr.find_all(["td","th"], recursive=False)]
            body_rows.append(cells)

        ncols = len(headers) if headers else (len(body_rows[0]) if body_rows else 0)
        nrows = len(body_rows) + (1 if headers else 0)
        if ncols <= 0 or nrows <= 0:
            self.warnings.append("Empty <table> skipped.")
            return 0.0

        # Height estimate
        row_h_in = max(0.3, self._line_height_in(tokens, body=True) + 0.05)
        total_h_in = nrows * row_h_in + 0.05

        shape = slide.shapes.add_table(nrows, ncols, Inches(x_in), Inches(y_in), Inches(w_in), Inches(total_h_in))
        tbl = shape.table

        # Column widths
        col_w_in = w_in / ncols
        for i in range(ncols):
            tbl.columns[i].width = Inches(col_w_in)

        r_offset = 0
        if headers:
            for j, text in enumerate(headers[:ncols]):
                cell = tbl.cell(0, j)
                tf = cell.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(0)
                p.line_spacing = 1.3
                run = p.add_run()
                self._set_font(run, name=tokens.font_sans, size_pt=tokens.fs_body_pt, color=(255,255,255), bold=True)
                run.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(tokens.accent)
                tbl.rows[0].height = Inches(row_h_in)
            r_offset = 1

        for i, row in enumerate(body_rows):
            for j in range(ncols):
                txt = row[j] if j < len(row) else ""
                cell = tbl.cell(i + r_offset, j)
                tf = cell.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(0)
                p.line_spacing = 1.3
                r = p.add_run()
                self._set_font(r, name=tokens.font_sans, size_pt=tokens.fs_body_pt, color=tokens.ink, bold=False)
                r.text = txt
            tbl.rows[i + r_offset].height = Inches(row_h_in)

        return total_h_in

    # ---------- Content block renderer ----------

    def _add_content_block(self, slide, tokens: DesignTokens, html_block, base_dir: Path, x_in: float, y_in: float, w_in: float, h_in: float):
        """Render an HTML <section class="content">...</section> using a simple flow layout."""
        self._base_dir = base_dir
        cursor_y = y_in
        max_y = y_in + h_in

        def advance(dh: float):
            nonlocal cursor_y
            cursor_y += dh + 0.12  # small gap after each block

        for child in html_block.children:
            if cursor_y >= max_y:
                break
            name = getattr(child, "name", None)
            if name is None:
                text = str(child).strip()
                if text:
                    used = self._add_text_block(slide, tokens, [text], x_in, cursor_y, w_in)
                    advance(used)
                continue

            if name in ("p", "div"):
                txt = child.get_text(" ", strip=True)
                if txt:
                    used = self._add_text_block(slide, tokens, [txt], x_in, cursor_y, w_in)
                    advance(used)

            elif name in ("h1", "h2", "h3"):
                txt = child.get_text(" ", strip=True)
                if txt:
                    used = self._add_text_block(slide, tokens, [txt], x_in, cursor_y, w_in, bold=True)
                    advance(used)

            elif name == "ul":
                items = []
                for li in child.find_all("li", recursive=False):
                    t = li.get_text(" ", strip=True)
                    if t:
                        items.append(t)
                if items:
                    used = self._add_list_block(slide, tokens, items, x_in, cursor_y, w_in)
                    advance(used)

            elif name == "img":
                src = child.get("src") or child.get("data-src") or ""
                used = self._add_image_block(slide, tokens, src, x_in, cursor_y, w_in)
                if used <= 0.0:
                    self.warnings.append("Skipped an <img> (missing or unreadable src).")
                advance(used if used>0 else 0.0)

            elif name == "figure":
                img = child.find("img")
                if img:
                    used = self._add_image_block(slide, tokens, img.get("src") or img.get("data-src") or "", x_in, cursor_y, w_in)
                    advance(used if used>0 else 0.0)
                    cap = child.find("figcaption")
                    if cap:
                        cap_txt = cap.get_text(" ", strip=True)
                        if cap_txt:
                            used2 = self._add_text_block(slide, tokens, [cap_txt], x_in, cursor_y, w_in, small=True)
                            advance(used2)
                else:
                    self.warnings.append("<figure> without <img> was ignored.")

            elif name == "table":
                used = self._add_table_block(slide, tokens, child, x_in, cursor_y, w_in)
                advance(used if used>0 else 0.0)

            else:
                t = child.get_text(" ", strip=True)
                if t:
                    used = self._add_text_block(slide, tokens, [t], x_in, cursor_y, w_in)
                    advance(used)

    def _add_regular_slide(self, prs: Presentation, tokens: DesignTokens, eyebrow: Optional[str], title: Optional[str],
                           content_sections: List, base_dir: Path, footer_left: str, page_text: str, two_col: bool = False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_header(slide, tokens, eyebrow, title)

        left_margin = 0.6
        right_margin = 0.6
        top_after_header = 0.4 + 1.45
        bottom_margin = 0.8
        width = self.page_width_in - left_margin - right_margin
        height = self.page_height_in - top_after_header - bottom_margin - 0.5

        if two_col and len(content_sections) >= 2:
            total_fr = 2.2
            left_w = width * (1.2 / total_fr)
            right_w = width * (1.0 / total_fr)
            gap_in = 0.3
            left_w -= gap_in/2
            right_w -= gap_in/2
            self._add_content_block(slide, tokens, content_sections[0], base_dir, left_margin, top_after_header + 0.2, left_w, height - 0.2)
            self._add_content_block(slide, tokens, content_sections[1], base_dir, left_margin + left_w + gap_in, top_after_header + 0.2, right_w, height - 0.2)
        else:
            if content_sections:
                self._add_content_block(slide, tokens, content_sections[0], base_dir, left_margin, top_after_header + 0.2, width, height - 0.2)

        self._add_footer(slide, tokens, footer_left, page_text)
        return slide

    def set_full_bg_image(self, slide, image_path: str):
        from pptx.util import Inches
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(self.page_width_in), height=Inches(self.page_height_in))

    def convert(self, html_path: str | Path, pptx_path: str | Path) -> Dict:
        html_path = Path(html_path)
        pptx_path = Path(pptx_path)
        self._base_dir = html_path.parent
        self.warnings = []

        html = html_path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")

        style_tag = soup.find("style")
        varmap = parse_css_vars(style_tag.get_text() if style_tag else "")
        tokens = DesignTokens.from_css_vars(varmap)

        slides_dom = soup.select("article.slide")
        total = len(slides_dom)

        prs = Presentation()
        prs.slide_width = Inches(self.page_width_in)
        prs.slide_height = Inches(self.page_height_in)

        for idx, s in enumerate(slides_dom, start=1):
            is_title = "title-slide" in (s.get("class") or [])
            footer_left = s.select_one(".footer div")
            footer_left_text = footer_left.get_text(" ", strip=True) if footer_left else ""
            page_text = f"{idx} / {total}"

            if is_title:
                kicker = s.select_one(".hero .kicker")
                t = s.select_one(".hero .title")
                sub = s.select_one(".hero .subtitle")
                self._add_title_slide(
                    prs, tokens,
                    kicker.get_text(" ", strip=True) if kicker else "",
                    t.get_text("\n", strip=True) if t else "",
                    sub.get_text(" ", strip=True) if sub else "",
                    footer_left_text, page_text
                )
            else:
                header = s.select_one(".header")
                eyebrow = header.select_one(".eyebrow").get_text(" ", strip=True) if header and header.select_one(".eyebrow") else None
                title = header.select_one(".title").get_text(" ", strip=True) if header and header.select_one(".title") else None

                contents = s.select(".content")
                two_col = any("two-col" in (c.get("class") or []) for c in contents)
                content_sections: List = []
                if contents:
                    c0 = contents[0]
                    if "two-col" in (c0.get("class") or []):
                        divs = c0.find_all("div", recursive=False)
                        if len(divs) >= 2:
                            content_sections = [divs[0], divs[1]]
                            two_col = True
                        else:
                            ul = c0.find("ul")
                            if ul:
                                lis = ul.find_all("li", recursive=False)
                                mid = len(lis)//2 or 1
                                left_ul = type(ul)(ul.name, ul.attrs.copy())
                                right_ul = type(ul)(ul.name, ul.attrs.copy())
                                for li in lis[:mid]:
                                    left_ul.append(li)
                                for li in lis[mid:]:
                                    right_ul.append(li)
                                content_sections = [left_ul, right_ul]
                                two_col = True
                    else:
                        content_sections = [c0]

                self._add_regular_slide(prs, tokens, eyebrow, title, content_sections, self._base_dir, footer_left_text, page_text, two_col=two_col)

        prs.save(str(pptx_path))

        meta = {
            "slides": total,
            "fonts_pt": {
                "title_xl": tokens.fs_title_xl_pt,
                "title": tokens.fs_title_pt,
                "body": tokens.fs_body_pt,
                "small": tokens.fs_small_pt,
            },
            "colors_rgb": {
                "brand": tokens.brand,
                "accent": tokens.accent,
                "bg": tokens.bg,
                "ink": tokens.ink,
                "muted": tokens.muted,
                "line": tokens.line,
            },
            "font_sans": tokens.font_sans,
            "warnings": self.warnings,
        }
        return meta


def convert_html_to_pptx(html_path: str | Path, pptx_path: str | Path) -> Dict:
    conv = HTMLtoPPTXConverter()
    return conv.convert(html_path, pptx_path)
